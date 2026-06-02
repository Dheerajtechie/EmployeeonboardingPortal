from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    slides_data = [
        ("Team Introduction", [
            "Team name: Tech Mavericks",
            "Frontend Dev: Manav",
            "Backend Dev: Rohan",
            "Database Engineer: Priya",
            "DevOps Engineer: Arjun",
            "AI Engineer: Sneha",
            "Docs & Presentation: Vikram"
        ]),
        ("Problem Statement", [
            "Day-1 chaos for new hires.",
            "No clear direction on what to do.",
            "Manual paper forms tracking.",
            "Missing assets and delayed IT provisioning.",
            "HR spends too much time answering repetitive questions."
        ]),
        ("Technology Stack", [
            "Frontend: React JS, Tailwind CSS, Axios",
            "Backend: FastAPI (Python), JWT Auth",
            "AI Service: Flask, Groq API, FAISS, Sentence Transformers",
            "Database: Oracle DB (cx_Oracle)",
            "Containerization: Docker, Docker Compose",
            "Orchestration: Kubernetes (K8s)",
            "CI/CD: Jenkins pipeline",
            "Version Control: GitHub"
        ]),
        ("System Architecture", [
            "Three independent microservices communicating via REST.",
            "React App runs in browser, communicates with Backend & AI.",
            "FastAPI backend manages core logic and connects to Oracle.",
            "Flask AI service handles semantic search over FAISS & LLM.",
            "Docker and Kubernetes orchestrate all 3 containers seamlessly."
        ]),
        ("Database Design", [
            "10 highly normalized Oracle tables.",
            "Core Entities: USERS, DEPARTMENTS, DOCUMENTS.",
            "Tasks: ONBOARDING_TASKS, TASK_ASSIGNMENTS.",
            "Assets: ASSETS, ASSET_ASSIGNMENTS.",
            "Trainings & Buddies: TRAININGS, BUDDIES, BUDDY_CHECKINS."
        ]),
        ("New Hire Journey Demo", [
            "Login with temporary credentials & set password.",
            "View personalized onboarding checklist.",
            "Upload ID and documents securely.",
            "Acknowledge receipt of IT assets.",
            "Access mandatory training resources."
        ]),
        ("HR & IT & Buddy Demo", [
            "HR Dashboard tracks all new hires' completion %.",
            "HR verifies or rejects uploaded documents with comments.",
            "IT Admin assigns hardware directly from inventory.",
            "Buddy logs weekly check-in notes for assigned new hires."
        ]),
        ("AI Chatbot Demo", [
            "RAG (Retrieval-Augmented Generation) Architecture.",
            "FAISS vector store holds policy documents and FAQs.",
            "System fetches the user's personal pending tasks.",
            "Result: 'You still need to submit PAN. To get a laptop, go to Floor 3.'"
        ]),
        ("Docker/K8s Demo", [
            "docker-compose up --build instantiates full stack locally.",
            "Persistent volumes ensure uploaded documents aren't lost.",
            "Kubernetes YAMLs separate Deployments and Services.",
            "ConfigMaps and Secrets handle environment variables safely."
        ]),
        ("Jenkins Pipeline", [
            "Triggered by GitHub Webhooks on push to main.",
            "Stage 1: Checkout Code.",
            "Stage 2: Install dependencies & run PyTest.",
            "Stage 3: Build Docker images for all 3 services.",
            "Stage 4: Deploy new image tags to Kubernetes Cluster."
        ]),
        ("Challenges & Learnings", [
            "Challenge: Handling multipart file uploads in FastAPI and Docker volumes.",
            "Learning: How cx_Oracle session pooling dramatically improves API speeds.",
            "Challenge: Fixing the task assignment mismatch (TA_ID vs ASSIGNMENT_ID).",
            "Learning: Integrating in-memory FAISS vs persistent vector stores."
        ]),
        ("Future Scope", [
            "Email and SMS alerts via Twilio/SendGrid integration.",
            "Mobile-responsive app or dedicated React Native mobile app.",
            "Digital e-signature integration for NDA and Code of Conduct.",
            "Learning Management System (LMS) API integration.",
            "Exit management and offboarding module."
        ])
    ]

    for title, points in slides_data:
        slide_layout = prs.slide_layouts[1] # Bullet slide
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for i, point in enumerate(points):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

    prs.save("Employee_Onboarding_Portal_Presentation.pptx")
    print("Created Employee_Onboarding_Portal_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
